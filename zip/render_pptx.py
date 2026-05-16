from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import os

BASE_DIR = Path(__file__).resolve().parent
# Save to parent directory to avoid permission issues in zip/
OUTPUT_PATH = BASE_DIR.parent / "ORBIT_App.pptx"
IMAGE_PATH = BASE_DIR.parent / "WhatsApp Image 2026-04-08 at 2.33.52 PM.jpeg"

screenshot_dir = os.environ.get("SCREENSHOT_DIR")
if screenshot_dir:
    screenshot_base = Path(screenshot_dir)
else:
    screenshot_base = BASE_DIR

slide_images = sorted(screenshot_base.glob('slide-*.png'))
if not slide_images:
    raise SystemExit("No slide images found. Run render_app.js first.")

prs = Presentation()

# Add first slide containing the local image
if IMAGE_PATH.exists():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Set the image as slide background
    background = slide.background
    fill = background.fill
    fill.user_picture(str(IMAGE_PATH))
else:
    print(f"Warning: First slide image not found at {IMAGE_PATH}")

# Add the rendered app slide images after the intro image slide
for img_path in slide_images:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(img_path), Inches(0), Inches(0), width=prs.slide_width)

try:
    prs.save(str(OUTPUT_PATH))
    print(f"Generated PPTX: {OUTPUT_PATH}")
except PermissionError:
    print(f"Error: Permission denied when saving {OUTPUT_PATH}. Please make sure the file is closed in PowerPoint.")

# Remove the temporary screenshot files if they were stored in a temp directory
if screenshot_dir and screenshot_base != BASE_DIR:
    try:
        for img_path in screenshot_base.glob('slide-*.png'):
            try:
                img_path.unlink()
            except Exception:
                pass
        screenshot_base.rmdir()
    except Exception:
        pass
